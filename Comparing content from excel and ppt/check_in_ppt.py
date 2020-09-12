from pptx import Presentation
import pandas as pd

prs = Presentation('sample.pptx')

content = ""
content_list=[]
for slide in prs.slides:
    for shape in slide.shapes:
        content=content+shape.text
        content=content+"\n"
        content_list.append(shape.text)
"""
for slide in prs.slides:
    for shape in slide.shapes:
        print(shape.name,shape.text)
    """
excel=pd.read_excel('input.xlsx')
master = dict(zip(excel['Column'],excel['Status'])) 
checked_status=[]
for i in master:
    if i.lower() in content.lower():
        #print(i,"Yes",master[i])
        checked_status.append('Y')
    else:
        #print(i,"No",master[i])
        checked_status.append('N')
excel['Checked_status']=checked_status

counter=0
for i,j in zip(excel.Status,excel.Checked_status):
    #print(i,j)
    if counter < excel.Column.shape[0]:
        if i==j:
            #print(excel.Column[counter],"is present")
            pass
        else:
            if i == 'Y':
                print(excel.Column[counter],"--It is missing in the PPT. Please check")
            else:
                print(excel.Column[counter],"--It is wrongly included in PPT and not supposed to be present. Please check")
    else:
        break
    counter+=1
